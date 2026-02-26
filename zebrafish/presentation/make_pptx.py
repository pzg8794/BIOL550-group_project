#!/usr/bin/env python3
from __future__ import annotations

from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Inches, Pt


def _set_run_font(run, name: str | None = None, size: int | None = None, bold: bool | None = None):
    if name is not None:
        run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.level = 0
        p.font.size = Pt(20)
    first = True
    for bullet in bullets:
        p = tf.add_paragraph() if (subtitle or not first) else tf.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)
        first = False


def add_diagram_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    top = Inches(1.7)
    left = Inches(0.6)
    w = Inches(2.9)
    h = Inches(1.0)
    gap = Inches(0.4)

    boxes = [
        ("Learn SRA tools", "prefetch, dump, gzip\nFastQC"),
        ("Validate manually", "single SRR tests\nlog everything"),
        ("Organize shared data", "sra_runs\nfastqc_out\nfastx_out"),
        ("Split run IDs", "10 SRRs / member\nreproducible lists"),
        ("Automate pipeline", "2 jobs in parallel\n1 SRR at a time"),
        ("Deliverables + analysis", "QC review\nalignment / DE"),
    ]

    shapes = []
    for i, (hdr, body) in enumerate(boxes):
        x = left + (w + gap) * (i % 3)
        y = top + (h + Inches(0.7)) * (i // 3)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
        shape.line.color.rgb = RGBColor(80, 80, 80)
        tf = shape.text_frame
        tf.clear()

        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = hdr
        _set_run_font(r0, size=18, bold=True)

        p1 = tf.add_paragraph()
        p1.text = body
        p1.level = 0
        p1.font.size = Pt(14)
        shapes.append(shape)

    def arrow(from_shape, to_shape):
        x1 = from_shape.left + from_shape.width
        y1 = from_shape.top + from_shape.height // 2
        x2 = to_shape.left
        y2 = to_shape.top + to_shape.height // 2
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = RGBColor(80, 80, 80)
        conn.line.width = Pt(2)

    arrow(shapes[0], shapes[1])
    arrow(shapes[1], shapes[2])
    arrow(shapes[3], shapes[4])
    arrow(shapes[4], shapes[5])

    # Vertical connectors
    def down_arrow(from_shape, to_shape):
        x1 = from_shape.left + from_shape.width // 2
        y1 = from_shape.top + from_shape.height
        x2 = to_shape.left + to_shape.width // 2
        y2 = to_shape.top
        conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
        conn.line.color.rgb = RGBColor(80, 80, 80)
        conn.line.width = Pt(2)

    down_arrow(shapes[2], shapes[3])


def add_two_column_bullets(prs: Presentation, title: str, left_title: str, left: list[str], right_title: str, right: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4.6), Inches(5.2))
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.6), Inches(5.2))

    for box, hdr, bullets in [(left_box, left_title, left), (right_box, right_title, right)]:
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = hdr
        p.font.bold = True
        p.font.size = Pt(22)
        p.level = 0
        for b in bullets:
            q = tf.add_paragraph()
            q.text = b
            q.level = 0
            q.font.size = Pt(18)


def add_code_examples_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(9.0), Inches(5.3))
    tf = tb.text_frame
    tf.clear()

    def add_section(header: str, code: str, note: str):
        p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        r = p.add_run()
        r.text = header
        _set_run_font(r, size=18, bold=True)

        q = tf.add_paragraph()
        q.text = code
        q.level = 0
        for run in q.runs:
            _set_run_font(run, name="Courier New", size=14)

        n = tf.add_paragraph()
        n.text = note
        n.level = 0
        n.font.size = Pt(14)

        spacer = tf.add_paragraph()
        spacer.text = ""

    add_section(
        "SRA download + dump (server install)",
        "/usr/local/bin/sra_3.0.0/prefetch.3.0.0 -O /tmp/sra SRR34002437\n"
        "/usr/local/bin/sra_3.0.0/fasterq-dump-orig.3.0.0 /tmp/sra/SRR34002437/SRR34002437.sra --split-files --threads 1 -O /tmp/fastq --temp /tmp/tmp\n"
        "gzip -f /tmp/fastq/SRR34002437_1.fastq /tmp/fastq/SRR34002437_2.fastq",
        "We use fasterq-dump when we want lower wall-clock time and can tolerate temporary uncompressed FASTQs; we keep threads low to avoid I/O contention.",
    )

    add_section(
        "FastQC (paired-end)",
        "/usr/local/bin/FASTQC_11.9/fastqc -t 1 -o /home/zebrafish/fastqc_out \\\n"
        "/home/zebrafish/sra_runs/SRR34002437_1.fastq.gz /home/zebrafish/sra_runs/SRR34002437_2.fastq.gz",
        "FastQC runs as a separate worker so QC overlaps with the next download. Outputs are HTML + ZIP per mate.",
    )

    add_section(
        "Trimming (FASTX example)",
        "mkdir -p /home/zebrafish/fastx_out\n"
        "zcat /home/zebrafish/sra_runs/SRR34002437_1.fastq.gz | fastx_trimmer -Q33 -f 1 -l 150 | gzip -c > /home/zebrafish/fastx_out/SRR34002437_1.trim.fastq.gz",
        "Trimming is optional and runs after raw QC; it produces separate trimmed outputs so raw FASTQs remain unchanged.",
    )


def main():
    prs = Presentation()
    today = date.today().strftime("%b %d, %Y")

    add_title_slide(
        prs,
        "Zebrafish RNA-seq: SRA Download + QC Automation",
        f"BIOL550 Group Update • PRJNA1277581 • {today}\nTeam: Piter Garcia • Samuel Kopelev • Nikhi",
    )

    add_diagram_slide(prs, "Workflow overview (engineering view)")

    add_bullets_slide(
        prs,
        "Approach",
        [
            "Learn the SRA Toolkit workflow (prefetch → dump → gzip) and validate on small tests",
            "Implement and test the workflow following the course guide conventions (paired-end, reproducible logs)",
            "Automate with a minimal pipeline (sequential runs, resumable, easy monitoring)",
        ],
    )

    add_two_column_bullets(
        prs,
        "Data management / organization",
        "Shared directories (Sequoia)",
        [
            "/home/zebrafish/sra_runs (flat SRR*_1/_2.fastq.gz)",
            "/home/zebrafish/fastqc_out (FastQC html/zip)",
            "/home/zebrafish/fastx_out (optional trimmed outputs)",
            "/home/zebrafish/split_run_ids (per-member SRR lists)",
        ],
        "Run allocation",
        [
            "Split the 30 SRRs evenly (10 per member)",
            "Each member validates the SRA tools on a small test before automation",
            "Keep outputs group-accessible; keep pipeline state/logs in user home",
        ],
    )

    add_code_examples_slide(prs, "Execution examples (under the wrapper)")

    add_bullets_slide(
        prs,
        "Why fasterq-dump vs fastq-dump",
        [
            "fasterq-dump: faster conversion and supports threading; good for reducing total wall-clock time on the server",
            "Tradeoff: writes large uncompressed FASTQs temporarily, so it needs more scratch space and disk I/O",
            "fastq-dump --gzip: lower temporary disk footprint (writes .fastq.gz directly) but often slower",
            "We choose based on constraints: server speed vs laptop disk limits, while keeping 1 SRR at a time",
        ],
    )

    add_bullets_slide(
        prs,
        "Pipeline design",
        [
            "Two concurrent jobs: (A) downloader for one SRR at a time, (B) FastQC worker for completed SRRs",
            "Sequential processing to avoid failures and simplify resume (skip-if-exists semantics)",
            "Low thread counts to reduce contention; throughput is typically limited by network + disk I/O",
            "Optional trimming step runs after raw QC and writes to fastx_out without overwriting raw data",
        ],
    )

    add_bullets_slide(
        prs,
        "Status (by meeting) + next steps",
        [
            "All assigned SRRs downloaded as paired FASTQ.gz into the shared run directory",
            "FastQC completed for all runs; QC summaries reviewed for obvious issues",
            "Optional: run trimming and re-QC if needed (fastx_out kept separate from raw)",
            "Begin downstream analysis: alignment/quantification and differential expression workflow",
        ],
    )

    out = "/Users/pitergarcia/DataScience/Semester5/BIOL550/group_project/zebrafish/presentation/BIOL550_Zebrafish_SRA_Pipeline_Update.pptx"
    prs.save(out)
    print(out)


if __name__ == "__main__":
    main()

